#!/usr/bin/env python3
"""Extract Gaodun MBA courseware into provider-neutral C1 staging candidates."""

from __future__ import annotations

import argparse
import csv
import hashlib
import json
import sys
import zipfile
from pathlib import Path
from typing import Any

import openpyxl
import pypdf
from docx import Document
from openpyxl.utils import get_column_letter
from pypdf import PdfReader
from pptx import Presentation


PRIORITY = {3: 0, 4: 1, 2: 2}


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for block in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def sanitize_unicode(text: str) -> str:
    return text.encode("utf-8", errors="replace").decode("utf-8")


def write_text(path: Path, text: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(sanitize_unicode(text).rstrip() + "\n", encoding="utf-8")


def extract_pdf(source: Path) -> tuple[str, dict[str, Any], str]:
    document = PdfReader(source)
    sections: list[str] = []
    nonempty_pages = 0
    character_count = 0
    for index, page in enumerate(document.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            nonempty_pages += 1
            character_count += len(text)
        sections.append(f"## Page {index}\n\n{text}")
    pages = len(document.pages)
    average = character_count / pages if pages else 0.0
    sparse = pages > 0 and (average < 50 or nonempty_pages < pages / 2)
    loss = "Text layer only; page markers retained; layout and images are not reconstructed."
    if sparse:
        loss += " The PDF is predominantly scanned or text-sparse and remains an OCR candidate."
    return (
        "\n\n".join(sections),
        {
            "pages": pages,
            "nonempty_pages": nonempty_pages,
            "text_characters": character_count,
            "text_characters_per_page": average,
            "ocr_candidate": sparse,
        },
        loss,
    )


def extract_docx(source: Path) -> tuple[str, dict[str, Any], str]:
    document = Document(source)
    sections: list[str] = []
    paragraph_count = 0
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            sections.append(text)
            paragraph_count += 1
    table_rows = 0
    for table_index, table in enumerate(document.tables, start=1):
        sections.append(f"## Table {table_index}")
        for row_index, row in enumerate(table.rows, start=1):
            values = [cell.text.strip().replace("\n", " / ") for cell in row.cells]
            sections.append(f"Row {row_index}: " + " | ".join(values))
            table_rows += 1
    return (
        "\n\n".join(sections),
        {
            "paragraphs": paragraph_count,
            "tables": len(document.tables),
            "table_rows": table_rows,
        },
        "Paragraphs and table cells retained; embedded images, drawings, and exact layout are not reconstructed.",
    )


def stringify_cell(value: Any) -> str:
    if value is None:
        return ""
    return str(value).replace("\r", " ").replace("\n", " / ").replace("|", "\\|")


def extract_xlsx(source: Path) -> tuple[str, dict[str, Any], str]:
    workbook = openpyxl.load_workbook(source, read_only=True, data_only=False)
    sections: list[str] = []
    nonempty_rows = 0
    formula_cells = 0
    sheet_names: list[str] = []
    for worksheet in workbook.worksheets:
        sheet_names.append(worksheet.title)
        sections.append(f"## Sheet: {worksheet.title}")
        for row_index, row in enumerate(worksheet.iter_rows(), start=1):
            values: list[str] = []
            for column_index, cell in enumerate(row, start=1):
                if cell.value is None:
                    continue
                value = stringify_cell(cell.value)
                if isinstance(cell.value, str) and cell.value.startswith("="):
                    formula_cells += 1
                values.append(f"{get_column_letter(column_index)}={value}")
            if values:
                nonempty_rows += 1
                sections.append(f"Row {row_index}: " + " | ".join(values))
    workbook.close()
    return (
        "\n\n".join(sections),
        {
            "sheets": sheet_names,
            "nonempty_rows": nonempty_rows,
            "formula_cells": formula_cells,
        },
        "Non-empty rows, sheet names, values, and formulas retained; styles, charts, and workbook interaction are not reconstructed.",
    )


def extract_pptx(source: Path) -> tuple[str, dict[str, Any], str]:
    presentation = Presentation(source)
    sections: list[str] = []
    text_shapes = 0
    for slide_index, slide in enumerate(presentation.slides, start=1):
        sections.append(f"## Slide {slide_index}")
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                sections.append(text)
                text_shapes += 1
    return (
        "\n\n".join(sections),
        {"slides": len(presentation.slides), "text_shapes": text_shapes},
        "Accessible text shapes retained; visual-only elements, animations, and exact slide layout are not reconstructed.",
    )


def extract_zip(source: Path) -> tuple[dict[str, Any], dict[str, Any], str]:
    with zipfile.ZipFile(source) as archive:
        entries = [
            {
                "path": item.filename,
                "size": item.file_size,
                "compressed_size": item.compress_size,
                "crc32": f"{item.CRC:08x}",
                "is_directory": item.is_dir(),
            }
            for item in archive.infolist()
        ]
    result = {"schema": "babata.courseware.archive-inventory/v1", "entries": entries}
    return (
        result,
        {"entries": len(entries), "files": sum(not item["is_directory"] for item in entries)},
        "Archive member names, sizes, and CRC values retained; member contents were not recursively interpreted.",
    )


def load_rows(path: Path) -> list[dict[str, str]]:
    with path.open("r", encoding="utf-8-sig", newline="") as stream:
        rows = list(csv.DictReader(stream))
    return [
        row
        for row in rows
        if row["module_type"] == "courseware" and row["coverage_status"] == "pending"
    ]


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--ledger", required=True, type=Path)
    parser.add_argument("--output-dir", required=True, type=Path)
    parser.add_argument("--limit", type=int)
    args = parser.parse_args()

    rows = load_rows(args.ledger)
    rows.sort(key=lambda row: (PRIORITY.get(int(row["course_order"]), 99), int(row["course_order"]), int(row["module_id"])))
    if args.limit is not None:
        rows = rows[: args.limit]

    results_dir = args.output_dir / "results" / "courseware"
    manifest_path = args.output_dir / "courseware-manifest.json"
    items: list[dict[str, Any]] = []
    previous_by_module: dict[int, dict[str, Any]] = {}
    if manifest_path.exists():
        previous = json.loads(manifest_path.read_text(encoding="utf-8"))
        previous_by_module = {int(item["module_id"]): item for item in previous.get("items", [])}
    extractors = {
        ".pdf": (extract_pdf, "extracted_text", "pypdf", pypdf.__version__),
        ".docx": (extract_docx, "extracted_text", "python-docx", "1.2.0"),
        ".xlsx": (extract_xlsx, "extracted_text", "openpyxl", openpyxl.__version__),
        ".pptx": (extract_pptx, "extracted_text", "python-pptx", "1.0.2"),
        ".zip": (extract_zip, "structured_result", "python-zipfile", f"{sys.version_info.major}.{sys.version_info.minor}"),
    }

    for position, row in enumerate(rows, start=1):
        module_id = row["module_id"]
        source = Path(row["target_path"])
        extension = row["extension"].lower()
        previous_item = previous_by_module.get(int(module_id))
        if previous_item and previous_item.get("status") in {"candidate_ready", "registered"}:
            output_file = previous_item.get("output_file")
            if output_file and Path(output_file).is_file() and sha256_file(Path(output_file)) == previous_item.get("output_sha256"):
                items.append(previous_item)
                print(f"[{position}/{len(rows)}] module:{module_id} reused {row['title']}", flush=True)
                continue
        item: dict[str, Any] = {
            "module_id": int(module_id),
            "course_order": int(row["course_order"]),
            "course": row["course"],
            "phase": row["phase"],
            "title": row["title"],
            "source_path": str(source),
            "item_id": row["item_id"],
            "revision_id": row["revision_id"],
            "asset_id": row["asset_id"],
            "asset_sha256": row["asset_sha256"],
            "extension": extension,
            "status": "pending",
            "registration": None,
        }
        try:
            actual_hash = sha256_file(source)
            if actual_hash != row["asset_sha256"]:
                raise ValueError(f"asset hash mismatch: expected {row['asset_sha256']}, got {actual_hash}")
            if extension == ".doc":
                item.update(status="needs_legacy_conversion", error="Legacy .doc requires an external converter.")
            elif extension not in extractors:
                item.update(status="unsupported", error=f"Unsupported extension: {extension}")
            else:
                extractor, kind, model, version = extractors[extension]
                content, details, loss_notes = extractor(source)
                suffix = ".json" if kind == "structured_result" else ".md"
                output = results_dir / f"module-{module_id}{suffix}"
                if kind == "structured_result":
                    output.parent.mkdir(parents=True, exist_ok=True)
                    write_text(output, json.dumps(content, ensure_ascii=False, indent=2))
                else:
                    heading = f"# {row['title']}\n\nCourse: {row['course']}\n\nPhase: {row['phase']}"
                    write_text(output, heading + "\n\n" + str(content))
                item.update(
                    status="candidate_ready",
                    kind=kind,
                    provider="local_extract",
                    model=model,
                    tool_version=str(version),
                    output_file=str(output),
                    output_sha256=sha256_file(output),
                    params={
                        "scope": "issue-149-gaodun-mba-c1",
                        "source_authority": "website",
                        "website_module_id": int(module_id),
                        "provider_input_sha256": actual_hash,
                        "preprocessing": ["Full deterministic local extraction; source C0 remained read-only."],
                        **details,
                    },
                    loss_notes=loss_notes,
                )
        except Exception as error:  # Keep the batch moving and preserve readable identity.
            item.update(status="failed", error=f"{type(error).__name__}: {error}")
        items.append(item)
        print(f"[{position}/{len(rows)}] module:{module_id} {item['status']} {row['title']}", flush=True)
        manifest_path.parent.mkdir(parents=True, exist_ok=True)
        write_text(
            manifest_path,
            json.dumps(
                {"schema": "babata.mba.courseware-staging/v1", "items": items},
                ensure_ascii=False,
                indent=2,
            ),
        )

    counts: dict[str, int] = {}
    for item in items:
        counts[item["status"]] = counts.get(item["status"], 0) + 1
    print(json.dumps({"total": len(items), "by_status": counts, "manifest": str(manifest_path)}, ensure_ascii=False))
    return 0 if counts.get("failed", 0) == 0 else 1


if __name__ == "__main__":
    raise SystemExit(main())
